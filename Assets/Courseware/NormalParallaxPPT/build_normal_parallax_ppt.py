from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT_DIR = Path(r"c:\Users\lafa\Shader Course URP\Assets\Courseware\NormalParallaxPPT")
IMG_DIR = OUT_DIR / "images"
OUT_FILE = OUT_DIR / "法线贴图与视差贴图_两课时课件_v3.pptx"


def set_bg(slide, color=(245, 248, 252)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = RGBColor(23, 43, 77)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.25), Inches(11.8), Inches(0.6))
        sp = sub_box.text_frame.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(18)
        sr.font.color.rgb = RGBColor(74, 85, 104)


def add_bullets(slide, bullets, x=0.8, y=1.9, w=6.6, h=5.0, font_size=22):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(10)
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(34, 45, 65)


def add_image(slide, img_path, x=7.1, y=1.9, w=5.8, h=5.0):
    if not img_path.exists():
        return
    slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), Inches(w), Inches(h))


def add_code_block(slide, code, x=0.8, y=3.6, w=12.0, h=2.8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 41, 55)
    line = box.line
    line.color.rgb = RGBColor(55, 65, 81)

    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = code
    r.font.name = "Consolas"
    r.font.size = Pt(14)
    r.font.color.rgb = RGBColor(229, 231, 235)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    s = prs.slides.add_slide(blank)
    set_bg(s, (239, 246, 255))
    add_title(s, "法线贴图与视差贴图实现原理", "两课时教学课件（URP 实战向 / 初学者友好）")
    add_bullets(
        s,
        [
            "适用对象：Unity/URP 初学者，具备基础 Shader 阅读能力",
            "课程目标：理解原理 -> 看懂代码 -> 能自己调参数和排错",
            "案例基础：Phong + Normal + Parallax + AO + SpecMask"
        ],
        y=2.1, w=11.5, h=3.0, font_size=24
    )

    # Slide 2
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "课程安排（2 课时）")
    add_bullets(
        s,
        [
            "第 1 课时（90 分钟）：法线贴图与 TBN 空间",
            "第 2 课时（90 分钟）：视差贴图与完整光照整合",
            "每课时结构：概念讲解 40% + 案例拆解 40% + 实操与答疑 20%"
        ],
        w=11.8
    )

    # Slide 3
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "学习目标")
    add_bullets(
        s,
        [
            "知道法线贴图为何能“看起来凹凸”，但不增加模型面数",
            "理解切线空间与 TBN：为何法线贴图可跨模型复用",
            "掌握视差贴图 UV 偏移原理与迭代思路",
            "能在 URP 中写出可维护的 Normal + Parallax Shader"
        ],
        w=11.8
    )

    # Slide 4
    s = prs.slides.add_slide(blank)
    set_bg(s, (247, 250, 252))
    add_title(s, "第1课时：法线贴图核心认知", "为什么“蓝紫色贴图”能骗过光照")
    add_bullets(
        s,
        [
            "核心思想：改法线，不改几何",
            "RGB <-> 法线向量映射：x/y/z 对应 r/g/b",
            "蓝通道偏高的原因：多数法线整体朝 +Z",
            "关键难点：法线贴图是切线空间，不是世界空间"
        ],
        w=6.6, font_size=21
    )
    add_image(s, IMG_DIR / "img_1.png")

    # Slide 5
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "法线向量编码与解码")
    add_bullets(
        s,
        [
            "编码：rgb = normal * 0.5 + 0.5",
            "解码：normal = rgb * 2.0 - 1.0",
            "法线贴图采样后必须 normalize()",
            "导入设置错误（sRGB/Normal Type）会导致“塑料感”"
        ],
        w=11.8
    )
    add_code_block(
        s,
        "float3 normalTS = UnpackNormal(SAMPLE_TEXTURE2D(_NormalMap, sampler_NormalMap, uv));\n"
        "normalTS.xy *= _NormalIntensity;\n"
        "float3 normalWS = normalize(mul(normalTS, tbn));"
    )

    # Slide 6
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "TBN 与切线空间（重点）")
    add_bullets(
        s,
        [
            "T (Tangent)：沿纹理 U 方向",
            "B (Bitangent)：沿纹理 V 方向",
            "N (Normal)：几何法线方向",
            "TBN 矩阵作用：把切线空间向量转到世界空间"
        ],
        w=6.5
    )
    # 更准确的 TBN/切线空间示意图（LearnOpenGL）
    add_image(s, IMG_DIR / "tbn_vectors.png")

    # Slide 7
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "第1课时案例拆解（URP）")
    add_bullets(
        s,
        [
            "顶点阶段输出：positionWS / normalWS / tangentWS / bitangentWS / uv",
            "片元阶段流程：采样法线 -> TBN 变换 -> 参与光照",
            "主光 + 附加光分离：便于排查每个光源贡献",
            "调试建议：先关 Spec 和 AO，只看 NdotL 是否正常"
        ],
        w=11.8
    )

    # Slide 8
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "第1课时实操任务")
    add_bullets(
        s,
        [
            "任务A：调 _NormalIntensity，观察砖缝深浅变化",
            "任务B：翻转法线 G 通道，观察光照方向错误",
            "任务C：关闭 normalize()，观察高光破碎现象",
            "课堂产出：一份可解释“为什么这么写”的笔记"
        ],
        w=11.8
    )

    # Slide 9
    s = prs.slides.add_slide(blank)
    set_bg(s, (247, 250, 252))
    add_title(s, "第2课时：视差贴图原理", "UV 偏移模拟高度，不是真位移")
    add_bullets(
        s,
        [
            "输入：高度图(Height/Parallax Map) + 视线方向 viewDirTS",
            "核心：沿 viewDirTS.xy 偏移 UV",
            "本质：采样坐标变化造成“观察角度变化感”",
            "特点：几何不变，轮廓不变，细节感增强"
        ],
        w=6.6, font_size=21
    )
    # 更准确的视差贴图原理图（LearnOpenGL）
    add_image(s, IMG_DIR / "parallax_principle.png")

    # Slide 10
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "视差偏移公式（基础版）")
    add_bullets(
        s,
        [
            "P = viewDirTS.xy * height * scale",
            "uv' = uv - P",
            "视角越斜，偏移越明显",
            "scale 过大时会出现纹理撕裂/游走"
        ],
        w=11.8
    )
    add_code_block(
        s,
        "float2 uvParallax = uv;\n"
        "for (int i = 0; i < 10; i++)\n"
        "{\n"
        "    float h = SAMPLE_TEXTURE2D(_ParallaxMap, sampler_ParallaxMap, uvParallax).r;\n"
        "    uvParallax -= (0.5 - h) * viewDirTS.xy * _Parallax * 0.01;\n"
        "}"
    )

    # Slide 11
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "为什么要迭代（Steep Parallax 思想）")
    add_bullets(
        s,
        [
            "单次偏移：快，但层次感不足",
            "多次迭代：更接近“沿深度走层”效果",
            "步数越高越真实，但像素成本更高",
            "工程做法：近景高步数，远景低步数"
        ],
        w=11.8
    )

    # Slide 12
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "法线 + 视差协同流程")
    add_bullets(
        s,
        [
            "先算 TBN，再得 viewDirTS",
            "先做视差得到 uvParallax",
            "再用 uvParallax 统一采样：Main/AO/Spec/Normal",
            "最后用 normalWS 进入光照模型（Phong 或 PBR）"
        ],
        w=11.8
    )

    # Slide 13
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "与你案例代码的对应关系（Phong.shader）")
    add_bullets(
        s,
        [
            "BuildTBN()：构造切线空间基底",
            "ApplyParallaxUVIterative()：视差 UV 迭代",
            "ApplyNormalTSToWS()：法线转世界空间",
            "EvaluatePhongLighting()：单灯漫反射+高光"
        ],
        w=11.8
    )

    # Slide 14
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "函数原理详解 1：BuildTBN()")
    add_bullets(
        s,
        [
            "作用：构造切线空间基底矩阵 TBN，把切线空间向量映射到世界空间",
            "输入：tangentWS / bitangentWS / normalWS（来自顶点阶段）",
            "实现关键：三者先 normalize，再按列组成 float3x3(T, B, N)",
            "数学意义：v_world = mul(v_tangent, TBN)",
            "常见坑：T/B/N 不正交或未归一化，会导致法线偏斜、镜面高光漂移"
        ],
        w=7.0, h=4.8, font_size=19
    )
    add_code_block(
        s,
        "float3x3 BuildTBN(float3 tangentWS, float3 bitangentWS, float3 normalWS)\n"
        "{\n"
        "    return float3x3(normalize(tangentWS), normalize(bitangentWS), normalize(normalWS));\n"
        "}",
        x=7.1, y=2.0, w=5.9, h=2.0
    )
    add_image(s, IMG_DIR / "tbn_shown.png", x=7.1, y=4.2, w=5.9, h=2.2)

    # Slide 15
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "函数原理详解 2：ApplyParallaxUVIterative()")
    add_bullets(
        s,
        [
            "作用：根据高度图和视线方向，迭代修正采样 UV，制造深度错觉",
            "核心变量：viewDirTS.xy（视线在切线平面的投影）",
            "每次迭代：读取 height，按 (0.5 - height) * viewDirTS.xy * scale 偏移 UV",
            "步数意义：steps 越大层次越平滑，但片元成本越高",
            "常见坑：scale 过大或视角太斜会产生拉伸/断裂"
        ],
        w=7.0, h=4.8, font_size=19
    )
    add_code_block(
        s,
        "for (int step = 0; step < steps; step++)\n"
        "{\n"
        "    float h = SAMPLE_TEXTURE2D(heightTex, samplerHeightTex, uvParallax).r;\n"
        "    uvParallax -= (0.5 - h) * viewDirTS.xy * parallaxScale * 0.01;\n"
        "}",
        x=7.1, y=2.0, w=5.9, h=2.3
    )
    add_image(s, IMG_DIR / "parallax_depth.png", x=7.1, y=4.4, w=5.9, h=2.0)

    # Slide 16
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "函数原理详解 3：ApplyNormalTSToWS()")
    add_bullets(
        s,
        [
            "作用：把法线贴图采样得到的 normalTS 转换到世界空间 normalWS",
            "normalTS.xy *= intensity：控制表面起伏“横向扰动”强弱",
            "转换：normalWS = normalize(mul(normalTS, TBN))",
            "为什么只放大 xy：z 是“朝外”分量，直接大改会破坏法线稳定性",
            "常见坑：遗漏 normalize 导致 NdotL/NdotH 异常、亮斑闪烁"
        ],
        w=7.0, h=4.8, font_size=19
    )
    add_code_block(
        s,
        "float3 ApplyNormalTSToWS(float3 normalTS, float3x3 tbn, float normalIntensity)\n"
        "{\n"
        "    normalTS.xy *= normalIntensity;\n"
        "    return normalize(mul(normalTS, tbn));\n"
        "}",
        x=7.1, y=2.0, w=5.9, h=2.2
    )
    add_image(s, IMG_DIR / "img_3.png", x=7.1, y=4.3, w=5.9, h=2.1)

    # Slide 17
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "函数原理详解 4：EvaluatePhongLighting()")
    add_bullets(
        s,
        [
            "作用：计算单光源贡献 = 漫反射 + 高光",
            "漫反射：diff = max(dot(N, L), 0) * atten * lightColor * baseColor",
            "高光：spec = pow(max(dot(N, H), 0), shininess) * atten * specIntensity * specMask",
            "halfDir H = normalize(L + V)：Blinn-Phong 稳定且成本低",
            "返回单灯结果后可在主流程累加主光和附加光"
        ],
        w=7.0, h=4.8, font_size=19
    )
    add_code_block(
        s,
        "float3 halfDir = normalize(lightDirWS + viewDirWS);\n"
        "float ndoth = max(0.0, dot(normalWS, halfDir));\n"
        "float3 specular = pow(ndoth, shininess) * diffTerm * lightColor * specIntensity * specMask;\n"
        "return diffuse + specular;",
        x=7.1, y=2.0, w=5.9, h=2.2
    )
    add_image(s, IMG_DIR / "img_4.png", x=7.1, y=4.3, w=5.9, h=2.1)

    # Slide 18
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "函数调用链总览（从输入到出图）")
    add_bullets(
        s,
        [
            "1) BuildTBN -> 得到空间变换基底",
            "2) viewDirWS -> viewDirTS -> ApplyParallaxUVIterative 得到 uvParallax",
            "3) 用 uvParallax 采样法线并 ApplyNormalTSToWS 得到 normalWS",
            "4) EvaluatePhongLighting 计算主光与附加光",
            "5) AO + ToneMapACES + Gamma 输出最终颜色"
        ],
        w=11.8
    )

    # Slide 19
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "常见问题与排错清单")
    add_bullets(
        s,
        [
            "法线方向反了：检查切线 handedness / 法线贴图导入类型",
            "视差抖动：减小 _Parallax，限制极斜视角偏移",
            "高光异常：确认 normalWS 与 light/view 均 normalize",
            "画面灰蒙：检查线性空间与手动 gamma 转换是否重复"
        ],
        w=11.8
    )

    # Slide 20
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "性能与质量平衡")
    add_bullets(
        s,
        [
            "优先级：法线贴图 > 视差贴图（性价比更高）",
            "移动端建议：降低视差步数，或仅关键材质启用",
            "可加 Keyword：_PARALLAX_OFF / _PARALLAX_LOW / _PARALLAX_HIGH",
            "LOD 策略：远处关闭视差，仅保留法线"
        ],
        w=11.8
    )

    # Slide 21
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title(s, "两课时练习安排")
    add_bullets(
        s,
        [
            "练习1：只开法线贴图，做“砖墙起伏感”",
            "练习2：加视差贴图，对比不同 _Parallax 和步数",
            "练习3：记录 FPS 与视觉质量，给出参数建议",
            "练习4：封装自己的 Shading Library 函数"
        ],
        w=11.8
    )

    # Slide 22
    s = prs.slides.add_slide(blank)
    set_bg(s, (243, 244, 246))
    add_title(s, "课后作业")
    add_bullets(
        s,
        [
            "作业A：在你的材质库中新增一个“可调视差级别”的 Shader",
            "作业B：写一页文档解释“为什么先视差后法线采样”",
            "作业C：提交一组对比图（关闭/开启法线与视差）"
        ],
        w=11.8
    )

    # Slide 23
    s = prs.slides.add_slide(blank)
    set_bg(s, (255, 251, 235))
    add_title(s, "图片与资料来源（网络）")
    add_bullets(
        s,
        [
            "Wikimedia Commons: File:Normal_map_example.png",
            "Wikimedia Commons: File:Normal_map_example - Map.png",
            "Wikimedia Commons: File:Normal_map_example - Result.png",
            "Wikimedia Commons: File:Normal_map_example with scene and result.png",
            "Wikimedia Commons: File:Normal_map_example - Scene.png",
            "LearnOpenGL: Normal Mapping（TBN/切线空间示意图）",
            "LearnOpenGL: Parallax Mapping（原理示意图与算法）"
        ],
        w=12.0, font_size=18
    )

    # Slide 24
    s = prs.slides.add_slide(blank)
    set_bg(s, (239, 246, 255))
    add_title(s, "总结")
    add_bullets(
        s,
        [
            "法线贴图：核心在“空间统一”和“法线正确变换”",
            "视差贴图：核心在“UV 偏移模型”和“步数/性能平衡”",
            "工程实践：函数库封装 + 调试可视化 + 参数分级"
        ],
        w=11.8
    )

    # Slide 25
    s = prs.slides.add_slide(blank)
    set_bg(s, (23, 43, 77))
    add_title(s, "Q&A", "欢迎继续提问：我可以再给你一版“讲师逐页讲稿”")
    # recolor title on dark background
    for shape in s.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor(255, 255, 255)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_FILE))
    print(f"Generated: {OUT_FILE}")


if __name__ == "__main__":
    build()
